from __future__ import annotations

import base64
import hashlib
import json
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import fitz
import google.genai as genai
import pdfplumber
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from google.genai import types as genai_types
from openai import OpenAI
from openpyxl import load_workbook
from pptx import Presentation
from sqlalchemy.orm import Session

from app.services.chroma_service import get_retrieval_collection
from app.services.embedding_service import (
    LOCAL_EMBEDDING_MODEL,
    LOCAL_EMBEDDING_PROVIDER,
    chunk_text,
    count_tokens,
    embed_texts_with_local_model,
)
from app.services.ingestion_service import SUPPORTED_DOCUMENT_SUFFIXES
from app.services.provider_settings_service import (
    ProviderSettingsError,
    resolve_chat_provider_config,
)
from app.storage.models import Chunk, Document, Source
from app.storage.repositories.processing_run_repository import ProcessingRunRepository


class ProcessingError(Exception):
    pass


PDF_VISION_MAX_PAGES = 80
PDF_VISION_TIMEOUT_SECONDS = 45
TEXT_DOCUMENT_SUFFIXES = frozenset({".txt", ".md", ".csv", ".json", ".html", ".htm", ".xml"})


@dataclass
class ExtractedSourceContent:
    title: str
    text: str
    source_url: str
    page_texts: list[str] | None = None


@dataclass
class VisionExtractionConfig:
    provider: str
    api_key: str
    model: str
    base_url: str | None = None


def process_sources(
    db: Session,
    project_id: uuid.UUID,
    job_id: uuid.UUID,
    sources: list[Source],
    chunk_size: int,
    chunk_overlap: int,
    parent_run_id: uuid.UUID | None = None,
) -> dict[str, int | str]:
    run_repo = ProcessingRunRepository(db)
    embedding_model = LOCAL_EMBEDDING_MODEL
    run_metadata = {
        "source_ids": [str(source.id) for source in sources],
        "source_count": len(sources),
        "chunk_size": chunk_size,
        "chunk_overlap": chunk_overlap,
        "provider": LOCAL_EMBEDDING_PROVIDER,
        "embedding_model": embedding_model,
    }
    config_hash = hashlib.sha256(
        json.dumps(
            {
                "chunk_size": chunk_size,
                "chunk_overlap": chunk_overlap,
                "embedding_model": embedding_model,
            },
            sort_keys=True,
        ).encode("utf-8")
    ).hexdigest()
    run = run_repo.create(
        project_id=project_id,
        job_id=job_id,
        run_type="processing",
        model_id=embedding_model,
        prompt_hash=None,
        config_hash=config_hash,
        retrieval_config=None,
        run_metadata=run_metadata,
        parent_run_id=parent_run_id,
    )

    documents_created = 0
    chunks_created = 0

    for source in sources:
        _replace_source_documents(db, source)
        extracted = _extract_source_content(db, source)
        document = Document(
            source_id=source.id,
            title=extracted.title,
            raw_path=source.storage_path,
            clean_text=extracted.text,
            token_count=count_tokens(extracted.text, model_name=embedding_model),
        )
        db.add(document)
        db.commit()
        db.refresh(document)
        documents_created += 1

        chunk_payloads = _build_chunk_payloads(
            source=source,
            document=document,
            text=extracted.text,
            source_url=extracted.source_url,
            page_texts=extracted.page_texts,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            embedding_model=embedding_model,
        )
        vectors = embed_texts_with_local_model(
            [item["content"] for item in chunk_payloads],
            model_name=embedding_model,
        )
        chunk_models: list[Chunk] = []
        for item in chunk_payloads:
            chunk_models.append(
                Chunk(
                    id=uuid.UUID(item["chunk_id"]),
                    document_id=document.id,
                    chunk_index=item["chunk_index"],
                    content=item["content"],
                    chroma_id=item["chroma_id"],
                    embedding_model=embedding_model,
                    chunk_metadata=item["metadata"],
                )
            )
        collection = get_retrieval_collection(embedding_model)
        collection.add(
            ids=[chunk.chroma_id for chunk in chunk_models if chunk.chroma_id],
            documents=[chunk.content for chunk in chunk_models],
            embeddings=vectors,
            metadatas=[chunk.chunk_metadata for chunk in chunk_models],
        )

        for chunk in chunk_models:
            db.add(chunk)
            chunks_created += 1

        db.commit()

    run = run_repo.update_metadata(
        run,
        {
            **run_metadata,
            "documents_created": documents_created,
            "chunks_created": chunks_created,
        },
    )

    return {
        "run_id": str(run.id),
        "documents_created": documents_created,
        "chunks_created": chunks_created,
    }


def _replace_source_documents(db: Session, source: Source) -> None:
    existing_documents = db.query(Document).filter(Document.source_id == source.id).all()
    if not existing_documents:
        return

    document_ids = [document.id for document in existing_documents]
    existing_chunks = db.query(Chunk).filter(Chunk.document_id.in_(document_ids)).all()
    chroma_ids = [chunk.chroma_id for chunk in existing_chunks if chunk.chroma_id]
    if chroma_ids:
        get_retrieval_collection(LOCAL_EMBEDDING_MODEL).delete(ids=chroma_ids)

    for chunk in existing_chunks:
        db.delete(chunk)
    for document in existing_documents:
        db.delete(document)
    db.commit()


def _extract_source_content(db: Session, source: Source) -> ExtractedSourceContent:
    if not source.storage_path:
        raise ProcessingError(f"Source '{source.id}' does not have a stored artifact.")

    storage_path = Path(source.storage_path)
    if not storage_path.exists():
        raise ProcessingError(f"Stored artifact '{storage_path}' does not exist.")

    if source.type == "file":
        return _extract_file_content(storage_path, _resolve_vision_config(db, source.project_id))
    if source.type in {"url", "arxiv", "google_drive", "notion", "slack", "confluence"}:
        return _extract_remote_payload(storage_path)
    raise ProcessingError(f"Unsupported source type '{source.type}'.")


def _extract_file_content(
    path: Path,
    vision_config: VisionExtractionConfig | None = None,
) -> ExtractedSourceContent:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_DOCUMENT_SUFFIXES:
        raise ProcessingError(f"Unsupported document format '{suffix or path.name}'.")

    try:
        if suffix == ".pdf":
            page_texts = _extract_pdf_pages(path, vision_config=vision_config)
            text = "\n".join(page.strip() for page in page_texts if page.strip())
            if not text.strip():
                raise ProcessingError(f"No readable text found in '{path.name}'.")
            return ExtractedSourceContent(
                title=path.name,
                text=text.strip(),
                source_url="",
                page_texts=page_texts,
            )
        if suffix == ".docx":
            text = _extract_docx_text(path)
        elif suffix == ".pptx":
            slide_texts = _extract_pptx_slides(path)
            text = "\n\n".join(slide.strip() for slide in slide_texts if slide.strip())
            if not text.strip():
                raise ProcessingError(f"No readable text found in '{path.name}'.")
            return ExtractedSourceContent(
                title=path.name,
                text=text.strip(),
                source_url="",
                page_texts=slide_texts,
            )
        elif suffix == ".xlsx":
            text = _extract_xlsx_text(path)
        elif suffix in TEXT_DOCUMENT_SUFFIXES:
            text = _extract_text_document(path, suffix=suffix)
        else:
            raise ProcessingError(f"Unsupported document format '{suffix}'.")
    except ProcessingError:
        raise
    except Exception as exc:
        raise ProcessingError(f"Failed to extract readable text from '{path.name}'.") from exc

    clean_text = text.strip()
    if not clean_text:
        raise ProcessingError(f"No readable text found in '{path.name}'.")

    return ExtractedSourceContent(title=path.name, text=clean_text, source_url="")


def _extract_text_document(path: Path, *, suffix: str) -> str:
    text = _decode_text_bytes(path.read_bytes())
    if suffix == ".json":
        try:
            return json.dumps(json.loads(text), ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            return text
    if suffix in {".html", ".htm"}:
        return _extract_html_text(text)
    return text


def _decode_text_bytes(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "cp1258", "cp1252"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _extract_html_text(raw_html: str) -> str:
    soup = BeautifulSoup(raw_html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    title = soup.title.get_text(" ", strip=True) if soup.title else ""
    body_text = soup.get_text("\n", strip=True)
    return "\n".join(part for part in [title, body_text] if part).strip()


def _extract_docx_text(path: Path) -> str:
    document = DocxDocument(path)
    parts: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                parts.append(" | ".join(values))

    return "\n".join(parts)


def _extract_pptx_slides(path: Path) -> list[str]:
    presentation = Presentation(path)
    slide_texts: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts = [f"Slide {slide_number}"]
        for shape in slide.shapes:
            parts.extend(_extract_pptx_shape_text(shape))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Notes: {notes}")
        slide_texts.append("\n".join(part for part in parts if part.strip()))

    return slide_texts


def _extract_pptx_shape_text(shape: Any) -> list[str]:
    parts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            parts.append(text)
    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                parts.append(" | ".join(values))
    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            parts.extend(_extract_pptx_shape_text(child_shape))
    return parts


def _extract_xlsx_text(path: Path) -> str:
    workbook = load_workbook(path, data_only=True, read_only=True)
    try:
        parts: list[str] = []
        for sheet in workbook.worksheets:
            parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [_format_spreadsheet_value(value) for value in row]
                values = [value for value in values if value]
                if values:
                    parts.append(" | ".join(values))
        return "\n".join(parts)
    finally:
        workbook.close()


def _format_spreadsheet_value(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def _extract_pdf_pages(
    path: Path,
    vision_config: VisionExtractionConfig | None = None,
) -> list[str]:
    fitz_pages: list[str] = []
    page_is_landscape: list[bool] = []

    with fitz.open(path) as document:
        for page in document:
            fitz_pages.append(page.get_text("text"))
            page_is_landscape.append(page.rect.width > page.rect.height)

    with pdfplumber.open(path) as pdf:
        plumber_pages = [(page.extract_text() or "") for page in pdf.pages]

    text_pages = _choose_pdf_text_pages(fitz_pages, plumber_pages)
    if vision_config is None or not _should_use_pdf_vision(text_pages, page_is_landscape):
        return text_pages

    return _augment_pdf_pages_with_vision(
        path,
        text_pages=text_pages,
        vision_config=vision_config,
    )


def _choose_pdf_text_pages(fitz_pages: list[str], plumber_pages: list[str]) -> list[str]:
    page_count = max(len(fitz_pages), len(plumber_pages))
    selected_pages: list[str] = []
    for index in range(page_count):
        fitz_text = fitz_pages[index].strip() if index < len(fitz_pages) else ""
        plumber_text = plumber_pages[index].strip() if index < len(plumber_pages) else ""
        selected_pages.append(_choose_better_page_text(fitz_text, plumber_text))
    return selected_pages


def _choose_better_page_text(fitz_text: str, plumber_text: str) -> str:
    if not fitz_text:
        return plumber_text
    if not plumber_text:
        return fitz_text
    if _pdf_text_quality_score(plumber_text) >= _pdf_text_quality_score(fitz_text):
        return plumber_text
    return fitz_text


def _pdf_text_quality_score(text: str) -> int:
    word_separators = text.count(" ") + text.count("\n") * 2
    replacement_penalty = text.count("\ufffd") * 20
    return len(text.strip()) + word_separators - replacement_penalty


def _resolve_vision_config(
    db: Session,
    project_id: uuid.UUID,
) -> VisionExtractionConfig | None:
    for provider in ("openai", "gemini"):
        try:
            config = resolve_chat_provider_config(db, project_id, provider)
        except ProviderSettingsError:
            continue
        return VisionExtractionConfig(
            provider=provider,
            api_key=config["api_key"],
            model=config["chat_model"],
            base_url=config.get("base_url"),
        )
    return None


def _should_use_pdf_vision(text_pages: list[str], page_is_landscape: list[bool]) -> bool:
    if not text_pages:
        return False
    non_empty_lengths = [len(page.strip()) for page in text_pages if page.strip()]
    if not non_empty_lengths:
        return True
    average_text_length = sum(non_empty_lengths) / len(non_empty_lengths)
    landscape_ratio = (
        sum(1 for is_landscape in page_is_landscape if is_landscape) / len(page_is_landscape)
        if page_is_landscape
        else 0
    )
    return average_text_length < 900 or landscape_ratio >= 0.6


def _augment_pdf_pages_with_vision(
    path: Path,
    *,
    text_pages: list[str],
    vision_config: VisionExtractionConfig,
) -> list[str]:
    augmented_pages: list[str] = []
    consecutive_vision_failures = 0
    with fitz.open(path) as document:
        for index, page in enumerate(document, start=1):
            text_layer = text_pages[index - 1] if index - 1 < len(text_pages) else ""
            if index > PDF_VISION_MAX_PAGES or consecutive_vision_failures >= 3:
                augmented_pages.append(text_layer)
                continue
            try:
                image_bytes = _render_pdf_page_png(page)
                visual_text = _describe_slide_image(
                    image_bytes=image_bytes,
                    page_number=index,
                    vision_config=vision_config,
                )
            except Exception:
                consecutive_vision_failures += 1
                visual_text = ""
            else:
                consecutive_vision_failures = 0
            augmented_pages.append(_merge_text_and_visual_summary(text_layer, visual_text))
    return augmented_pages


def _render_pdf_page_png(page: fitz.Page) -> bytes:
    pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
    return pixmap.tobytes("png")


def _describe_slide_image(
    *,
    image_bytes: bytes,
    page_number: int,
    vision_config: VisionExtractionConfig,
) -> str:
    prompt = (
        "You are extracting study notes from a Vietnamese university lecture slide. "
        "Read the slide image carefully. Return plain text in Vietnamese. Include: "
        "1) all visible headings and bullet text, 2) important text inside images, "
        "diagrams, tables, charts, and handwritten/scan regions when readable, "
        "3) a short explanation of what the visual content contributes. "
        f"This is slide/page {page_number}. Do not invent content that is not visible."
    )
    if vision_config.provider == "gemini":
        return _describe_slide_image_with_gemini(
            prompt=prompt,
            image_bytes=image_bytes,
            vision_config=vision_config,
        )
    return _describe_slide_image_with_openai(
        prompt=prompt,
        image_bytes=image_bytes,
        vision_config=vision_config,
    )


def _describe_slide_image_with_openai(
    *,
    prompt: str,
    image_bytes: bytes,
    vision_config: VisionExtractionConfig,
) -> str:
    kwargs: dict[str, Any] = {
        "api_key": vision_config.api_key,
        "max_retries": 0,
        "timeout": PDF_VISION_TIMEOUT_SECONDS,
    }
    if vision_config.base_url:
        kwargs["base_url"] = vision_config.base_url
    client = OpenAI(**kwargs)
    encoded_image = base64.b64encode(image_bytes).decode("ascii")
    response = client.chat.completions.create(
        model=vision_config.model,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{encoded_image}"},
                    },
                ],
            }
        ],
        temperature=0.1,
    )
    content = response.choices[0].message.content
    return content.strip() if content else ""


def _describe_slide_image_with_gemini(
    *,
    prompt: str,
    image_bytes: bytes,
    vision_config: VisionExtractionConfig,
) -> str:
    client = genai.Client(
        api_key=vision_config.api_key,
        http_options=genai_types.HttpOptions(timeout=PDF_VISION_TIMEOUT_SECONDS * 1000),
    )
    try:
        response = client.models.generate_content(
            model=vision_config.model,
            contents=[
                prompt,
                genai_types.Part.from_bytes(data=image_bytes, mime_type="image/png"),
            ],
        )
        content = getattr(response, "text", None)
        return content.strip() if content else ""
    finally:
        client.close()


def _merge_text_and_visual_summary(text_layer: str, visual_text: str) -> str:
    text_layer = text_layer.strip()
    visual_text = visual_text.strip()
    if text_layer and visual_text:
        return f"{text_layer}\n\nPhân tích nội dung trực quan của slide:\n{visual_text}"
    return visual_text or text_layer


def _extract_remote_payload(path: Path) -> ExtractedSourceContent:
    payload = json.loads(path.read_text(encoding="utf-8"))
    title = str(payload.get("title") or path.name)
    content = str(payload.get("content") or "").strip()
    source_url = str(payload.get("url") or payload.get("ingested_from") or "")
    if not content:
        raise ProcessingError(f"No readable content found in '{path.name}'.")
    return ExtractedSourceContent(title=title, text=content, source_url=source_url)


def _build_chunk_payloads(
    source: Source,
    document: Document,
    text: str,
    source_url: str,
    page_texts: list[str] | None,
    chunk_size: int,
    chunk_overlap: int,
    embedding_model: str,
) -> list[dict[str, Any]]:
    if page_texts:
        return _build_pdf_chunk_payloads(
            source=source,
            document=document,
            page_texts=page_texts,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            embedding_model=embedding_model,
        )

    chunks = chunk_text(
        text,
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        model_name=embedding_model,
    )
    payloads: list[dict[str, Any]] = []

    for index, chunk in enumerate(chunks):
        chunk_uuid = str(uuid.uuid4())
        retrieval_metadata = _build_retrieval_metadata(
            source=source,
            source_url=source_url,
            title=document.title or "",
        )
        payloads.append(
            {
                "chunk_id": chunk_uuid,
                "chroma_id": chunk_uuid,
                "chunk_index": index,
                "content": chunk,
                "metadata": {
                    "project_id": str(source.project_id),
                    "source_id": str(source.id),
                    "document_id": str(document.id),
                    "chunk_id": chunk_uuid,
                    "chunk_index": index,
                    "source_type": source.type,
                    "title": document.title or "",
                    "url": source_url,
                    **retrieval_metadata,
                },
            }
        )

    return payloads


def _build_pdf_chunk_payloads(
    *,
    source: Source,
    document: Document,
    page_texts: list[str],
    chunk_size: int,
    chunk_overlap: int,
    embedding_model: str,
) -> list[dict[str, Any]]:
    payloads: list[dict[str, Any]] = []
    chunk_index = 0

    for page_number, page_text in enumerate(page_texts, start=1):
        if not page_text.strip():
            continue
        page_chunks = chunk_text(
            page_text,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            model_name=embedding_model,
        )
        for chunk in page_chunks:
            chunk_uuid = str(uuid.uuid4())
            retrieval_metadata = _build_retrieval_metadata(
                source=source,
                source_url="",
                title=document.title or "",
            )
            payloads.append(
                {
                    "chunk_id": chunk_uuid,
                    "chroma_id": chunk_uuid,
                    "chunk_index": chunk_index,
                    "content": chunk,
                    "metadata": {
                        "project_id": str(source.project_id),
                        "source_id": str(source.id),
                        "document_id": str(document.id),
                        "chunk_id": chunk_uuid,
                        "chunk_index": chunk_index,
                        "source_type": source.type,
                        "title": document.title or "",
                        "url": "",
                        "page_number": page_number,
                        **retrieval_metadata,
                    },
                }
            )
            chunk_index += 1

    return payloads


def _build_retrieval_metadata(
    *,
    source: Source,
    source_url: str,
    title: str,
) -> dict[str, str | float]:
    source_metadata = source.source_metadata if isinstance(source.source_metadata, dict) else {}
    filters = source_metadata.get("retrieval_filters")
    filters = filters if isinstance(filters, dict) else {}
    quality = source_metadata.get("source_quality")
    quality = quality if isinstance(quality, dict) else {}

    tags = filters.get("tags")
    tag_values = [str(tag).strip().lower() for tag in tags if str(tag).strip()] if isinstance(tags, list) else []
    metadata: dict[str, str | float] = {
        "source_title": str(source_metadata.get("title") or title or source.original_uri or ""),
        "source_url": str(source_metadata.get("external_url") or source_url or ""),
    }
    for key in ("author", "published_at", "language"):
        value = filters.get(key)
        if value:
            metadata[key] = str(value)
    if tag_values:
        metadata["tags"] = ",".join(tag_values)
    for key in ("freshness_score", "trust_score", "ocr_confidence"):
        value = quality.get(key)
        if isinstance(value, int | float):
            metadata[key] = float(value)
    return metadata
