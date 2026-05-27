from openpyxl import Workbook
from pptx import Presentation

from app.services import processing_service
from app.storage.models import Chunk, Document, Job, ProcessingRun, Project, Source


class FakeCollection:
    def __init__(self):
        self.added_ids: list[str] = []
        self.deleted_ids: list[str] = []

    def add(self, ids, documents, embeddings, metadatas):
        self.added_ids.extend(ids)
        assert len(ids) == len(documents) == len(embeddings) == len(metadatas)

    def delete(self, ids):
        self.deleted_ids.extend(ids)


def test_process_sources_persists_run_and_replaces_existing_chunks(
    db_session,
    monkeypatch,
    tmp_path,
):
    project = Project(name="Phase 2", description="processing")
    db_session.add(project)
    db_session.commit()
    db_session.refresh(project)

    source_path = tmp_path / "notes.txt"
    source_path.write_text(
        "Alpha beta gamma delta epsilon zeta eta theta iota kappa lambda mu.",
        encoding="utf-8",
    )

    source = Source(
        project_id=project.id,
        type="file",
        original_uri="notes.txt",
        storage_path=str(source_path),
        checksum="checksum-source",
        status="completed",
        source_metadata={
            "source_quality": {
                "freshness_score": 1.0,
                "trust_score": 0.7,
                "ocr_confidence": None,
            },
            "retrieval_filters": {
                "author": "Research Team",
                "published_at": "2026-04-01",
                "language": "en",
                "tags": ["upload", "notes"],
            },
        },
    )
    db_session.add(source)
    db_session.commit()
    db_session.refresh(source)

    fake_collection = FakeCollection()
    monkeypatch.setattr(
        processing_service,
        "get_retrieval_collection",
        lambda embedding_model: fake_collection,
    )
    monkeypatch.setattr(
        processing_service,
        "embed_texts_with_local_model",
        lambda texts, model_name=None: [[float(index + 1)] * 3 for index, _ in enumerate(texts)],
    )

    first_job = Job(
        project_id=project.id,
        source_id=source.id,
        job_type="processing",
        status="running",
        progress=0,
    )
    db_session.add(first_job)
    db_session.commit()
    db_session.refresh(first_job)

    first_result = processing_service.process_sources(
        db=db_session,
        project_id=project.id,
        job_id=first_job.id,
        sources=[source],
        chunk_size=12,
        chunk_overlap=2,
    )

    first_chunk_ids = [
        chunk.chroma_id for chunk in db_session.query(Chunk).all() if chunk.chroma_id
    ]

    second_job = Job(
        project_id=project.id,
        source_id=source.id,
        job_type="processing",
        status="running",
        progress=0,
    )
    db_session.add(second_job)
    db_session.commit()
    db_session.refresh(second_job)

    second_result = processing_service.process_sources(
        db=db_session,
        project_id=project.id,
        job_id=second_job.id,
        sources=[source],
        chunk_size=12,
        chunk_overlap=2,
    )

    assert first_result["documents_created"] == 1
    assert second_result["documents_created"] == 1
    assert first_result["run_id"] != second_result["run_id"]
    assert db_session.query(Document).count() == 1
    assert db_session.query(Chunk).count() > 0
    assert fake_collection.deleted_ids == first_chunk_ids
    chunk_metadata = db_session.query(Chunk).first().chunk_metadata
    assert chunk_metadata["author"] == "Research Team"
    assert chunk_metadata["language"] == "en"
    assert chunk_metadata["tags"] == "upload,notes"
    assert chunk_metadata["trust_score"] == 0.7

    runs = db_session.query(ProcessingRun).order_by(ProcessingRun.created_at.asc()).all()
    assert len(runs) == 2
    assert runs[-1].run_metadata["documents_created"] == 1
    assert runs[-1].run_metadata["chunks_created"] == db_session.query(Chunk).count()


def test_process_pdf_sources_store_page_number_metadata(
    db_session,
    monkeypatch,
    tmp_path,
):
    project = Project(name="PDF pages", description="processing")
    db_session.add(project)
    db_session.commit()
    db_session.refresh(project)

    source_path = tmp_path / "paper.pdf"
    source_path.write_text("placeholder", encoding="utf-8")

    source = Source(
        project_id=project.id,
        type="file",
        original_uri="paper.pdf",
        storage_path=str(source_path),
        checksum="checksum-pdf",
        status="completed",
    )
    db_session.add(source)
    db_session.commit()
    db_session.refresh(source)

    fake_collection = FakeCollection()
    monkeypatch.setattr(
        processing_service,
        "get_retrieval_collection",
        lambda embedding_model: fake_collection,
    )
    monkeypatch.setattr(
        processing_service,
        "embed_texts_with_local_model",
        lambda texts, model_name=None: [[float(index + 1)] * 3 for index, _ in enumerate(texts)],
    )
    monkeypatch.setattr(
        processing_service,
        "_extract_pdf_pages",
        lambda path, vision_config=None: ["Page one text", "Page two text"],
    )

    job = Job(
        project_id=project.id,
        source_id=source.id,
        job_type="processing",
        status="running",
        progress=0,
    )
    db_session.add(job)
    db_session.commit()
    db_session.refresh(job)

    processing_service.process_sources(
        db=db_session,
        project_id=project.id,
        job_id=job.id,
        sources=[source],
        chunk_size=64,
        chunk_overlap=0,
    )

    page_numbers = {
        chunk.chunk_metadata["page_number"] for chunk in db_session.query(Chunk).all()
    }
    assert page_numbers == {1, 2}


def test_extract_pdf_pages_uses_vision_for_slide_pdf(monkeypatch, tmp_path):
    source_path = tmp_path / "slides.pdf"
    document = processing_service.fitz.open()
    page = document.new_page(width=1280, height=720)
    page.draw_rect(processing_service.fitz.Rect(0, 0, 1280, 720), color=(1, 1, 0))
    document.save(source_path)
    document.close()

    monkeypatch.setattr(
        processing_service,
        "_describe_slide_image",
        lambda image_bytes, page_number, vision_config: (
            "CHƯƠNG 1\nKhái niệm, đối tượng, phương pháp nghiên cứu."
        ),
    )

    page_texts = processing_service._extract_pdf_pages(
        source_path,
        vision_config=processing_service.VisionExtractionConfig(
            provider="openai",
            api_key="test-key",
            model="gpt-vision-test",
        ),
    )

    assert len(page_texts) == 1
    assert "CHƯƠNG 1" in page_texts[0]
    assert "phương pháp nghiên cứu" in page_texts[0]


def test_extract_pptx_content_reads_slide_text(tmp_path):
    source_path = tmp_path / "lecture.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Lecture 1"
    slide.placeholders[1].text = "Research question\nMethod summary"
    presentation.save(source_path)

    extracted = processing_service._extract_file_content(source_path)

    assert "Lecture 1" in extracted.text
    assert "Research question" in extracted.text
    assert extracted.page_texts is not None
    assert extracted.page_texts[0].startswith("Slide 1")


def test_extract_xlsx_content_reads_sheet_rows(tmp_path):
    source_path = tmp_path / "dataset.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Metrics"
    worksheet.append(["Name", "Score"])
    worksheet.append(["Alpha", 42])
    workbook.save(source_path)

    extracted = processing_service._extract_file_content(source_path)

    assert "Sheet: Metrics" in extracted.text
    assert "Name | Score" in extracted.text
    assert "Alpha | 42" in extracted.text


def test_choose_pdf_text_pages_prefers_better_spacing():
    pages = processing_service._choose_pdf_text_pages(
        ["Tư tưởng HồChí Minh là hệthống quan điểm."],
        ["Tư tưởng Hồ Chí Minh là hệ thống quan điểm."],
    )

    assert pages == ["Tư tưởng Hồ Chí Minh là hệ thống quan điểm."]
